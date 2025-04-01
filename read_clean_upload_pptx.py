from pptx import Presentation
import os
import json
from dotenv import load_dotenv
import re
from sharepoint_funcs_pptx_data import download_pptx
load_dotenv()




def clean_text(text):
    # Then, replace unwanted whitespace while keeping the regular space intact
    formatted_text = re.sub(r'[^\S ]+', '', text)
    return formatted_text


def clean_table_dimensions(table):
    """
    Returns a table with corrected dimensions.

    Tables in pptx sometimes do not follow correct dimensions for tables. For ex, the first row can have 2 cells
    and the second row can have 3 cells. This causes an error in reading the tables and often pptx reads the same data 
    twice to get dimension consistenncy. This function will add a "placeholder cell" to ensure that columns/cells dont 
    get repeated
    """
    for row in table.rows:
        print(len(row.cells))

    for col in table.columns:
        print(len(col.cells))

def extract_slide_tables(slide):
    """Return a list of table data dictionaries found in the slide.
    
    Each table dictionary includes:
      - table_no: sequential table index for the slide
      - columns: list of texts in the first row (if available)
      - rows: list of remaining rows (each a list of cell texts)

    Assumptions-> Assuming that each table has an assignned column heading. If not, the content of the first row will be 
    taken as heading.
    
    Params:
    slide -> slide object of pptx
      
    """

    
    tables = []
    table_index = 1
    for shape in slide.shapes:
        if shape.has_table:
            tbl = shape.table
            rows_data = []
            for row in tbl.rows:
                row_data = []
                for cell in row.cells:
                    cleaned_cell_text = clean_text(cell.text)
                    if cleaned_cell_text != "":
                        row_data.append(cleaned_cell_text)
                rows_data.append(row_data)
            # Assume first row represents column headers
            if rows_data:
                table_dict = {
                    "table_no": table_index,
                    "columns": rows_data[0],
                    "rows": rows_data[1:]
                }
            else:
                table_dict = {
                    "table_no": table_index,
                    "columns": [],
                    "rows": []
                }
            tables.append(table_dict)
            table_index += 1
    return tables

def pptx_to_json(ACCESS_TOKEN,file_path):
    """
    Return json string of pptx data. 

    Params:
    ACCESS_TOKEN -> To access sharepoint data
    file_path -> path of the pptx file relative to "Documents" document library in sharepoint
    """

    json_pptx_data = []
    pptx_file = download_pptx(ACCESS_TOKEN,file_path)

    if pptx_file == "Invalid file path":
        return "Invalid file path"
    elif pptx_file:
        pptx_file.seek(0)
        prs = Presentation(pptx_file)
       
        for slide in prs.slides:
            if slide.shapes.title:

                title = slide.shapes.title.text
            else:
                title = None
            slide_data = {
                "Title": clean_text(title),
                "Content":{"Tables":extract_slide_tables(slide)}
            }

            json_pptx_data.append(slide_data)
    
    else:
        return "Unable to get file data"


    return json.dumps(json_pptx_data,ensure_ascii=False)

